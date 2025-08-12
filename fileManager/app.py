from flask import Flask, request, render_template, send_from_directory, redirect, url_for, flash, jsonify
def get_dir_tree(path):
    tree = {'name': os.path.basename(path), 'path': path, 'type': 'folder', 'children': []}
    if os.path.isdir(path):
        try:
            for entry in os.listdir(path):
                full_path = os.path.join(path, entry)
                if os.path.isdir(full_path):
                    tree['children'].append(get_dir_tree(full_path))
        except Exception:
            pass
    return tree
import os
import shutil
from werkzeug.utils import secure_filename
from datetime import datetime
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pandas as pd

app = Flask(__name__)
app.secret_key = 'secret'  # Needed for flash messages

ALLOWED_EXTENSIONS = {
    'image': ['.png', '.jpg', '.jpeg', '.gif'],
    'pdf': ['.pdf'],
    'text': ['.txt', '.log'],
    'word': ['.doc', '.docx'],
    'excel': ['.xls', '.xlsx'],
    'ppt': ['.ppt', '.pptx'],
    'video': ['.mp4', '.webm']
}

@app.template_filter('datetimeformat')
def datetimeformat(value):
    return datetime.fromtimestamp(value).strftime('%Y-%m-%d %H:%M:%S')

def detect_filetype(filename):
    ext = os.path.splitext(filename)[1].lower()
    print(f"[DEBUG] detect_filetype: filename={filename}, ext={ext}")
    office_exts = ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']
    if ext in office_exts:
        print("[DEBUG] Returning 'office' for office file")
        return 'office'
    for filetype, extensions in ALLOWED_EXTENSIONS.items():
        if ext in extensions:
            print(f"[DEBUG] matched filetype={filetype}")
            return filetype
    print("[DEBUG] Returning 'other'")
    return 'other'

def read_office_file(full_path, filetype):
    print(f"[DEBUG] read_office_file: full_path={full_path}, filetype={filetype}")
    try:
        ext = os.path.splitext(full_path)[1].lower()
        if ext in ['.doc', '.docx']:
            print("[DEBUG] Reading DOC/DOCX file")
            doc = Document(full_path)
            return '\n'.join([para.text for para in doc.paragraphs])
        elif ext in ['.xls', '.xlsx']:
            print("[DEBUG] Reading XLS/XLSX file with pandas.read_excel")
            df = pd.read_excel(full_path)
            print(f"[DEBUG] DataFrame shape: {df.shape}")
            return df.to_html(classes='table table-bordered table-sm', index=False)
        elif ext in ['.ppt', '.pptx']:
            print("[DEBUG] Reading PPT/PPTX file")
            prs = Presentation(full_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return '\n'.join(texts)
    except Exception as e:
        print(f"[DEBUG] Exception in read_office_file: {e}")
        return f"Error reading office file: {e}"
    print("[DEBUG] Returning None from read_office_file")
    return None

@app.route('/', methods=['GET', 'POST'])
def index():
    directory = request.form.get('directory') or request.args.get('dir') or ''
    selected_file = request.args.get('file')
    selected_folder = request.args.get('folder')
    files = []
    tree = None

    if directory and os.path.exists(directory):
        tree = get_dir_tree(directory)
        if os.path.isdir(directory):
            files = sorted([f for f in os.listdir(directory) if os.path.isfile(os.path.join(directory, f))])
    else:
        directory = ''

    folder_files = []
    if selected_folder and os.path.isdir(selected_folder):
        try:
            folder_files = [f for f in os.listdir(selected_folder) if os.path.isfile(os.path.join(selected_folder, f))]
        except Exception:
            pass

    file_url = filetype = content = prev_file = next_file = file_stat = None

    if selected_file and directory and selected_file in files:
        full_path = os.path.join(directory, selected_file)
        if os.path.isfile(full_path):
            filetype = detect_filetype(selected_file)
            file_url = url_for('serve_file', dir=directory, filename=selected_file)
            print("james",filetype)
            if filetype in ['text', 'html']:
                with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            elif filetype == 'office':
                content = read_office_file(full_path, filetype)

            file_stat = os.stat(full_path)

            index = files.index(selected_file)
            prev_file = files[index - 1] if index > 0 else None
            next_file = files[index + 1] if index < len(files) - 1 else None

    return render_template('index.html', directory=directory, files=files, file=selected_file,
                           file_url=file_url, filetype=filetype, content=content,
                           prev_file=prev_file, next_file=next_file, file_stat=file_stat, os=os,
                           tree=tree, selected_folder=selected_folder, folder_files=folder_files)
@app.route('/get_tree', methods=['POST'])
def get_tree_route():
    path = request.form.get('path')
    if path and os.path.exists(path):
        tree = get_dir_tree(path)
        return jsonify(tree)
    return jsonify({'error': 'Path not found'})

@app.route('/list_folder_files', methods=['POST'])
def list_folder_files():
    folder_path = request.form.get('folder_path')
    files = []
    if folder_path and os.path.isdir(folder_path):
        try:
            files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
        except Exception:
            pass
    return jsonify({'files': files, 'folder': folder_path})

@app.route('/view')
def view_file():
    return redirect(url_for('index', dir=request.args.get('dir'), file=request.args.get('file')))

@app.route('/delete', methods=['POST'])
def delete_file():
    directory = request.args.get('dir')
    filename = request.args.get('file')

    if directory and filename:
        full_path = os.path.join(directory, filename)
        deleted_folder = os.path.join(directory, "deleted")
        os.makedirs(deleted_folder, exist_ok=True)

        deleted_path = os.path.join(deleted_folder, filename)
        try:
            if os.path.exists(full_path):
                shutil.move(full_path, deleted_path)
        except Exception as e:
            flash(f"Error deleting file: {e}")
            return redirect(url_for('index', dir=directory, file=filename))

        remaining_files = sorted([f for f in os.listdir(directory) if os.path.isfile(os.path.join(directory, f))])
        if remaining_files:
            next_file = remaining_files[0]
            return redirect(url_for('index', dir=directory, file=next_file))

    return redirect(url_for('index', dir=directory))

@app.route('/file')
def serve_file():
    directory = request.args.get('dir')
    filename = request.args.get('filename')
    if directory and filename:
        return send_from_directory(directory, filename)
    return '', 404

if __name__ == '__main__':
    app.run(debug=True)
